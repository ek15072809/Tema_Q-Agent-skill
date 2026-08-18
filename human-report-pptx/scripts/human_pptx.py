"""human_pptx.py — "人が作った" 報告・研究発表スタイルのスライドビルダー。

pptx_skill.py と同じ考え方（1関数=1スライド種、共通パレット/フォント、
python-pptx の上に薄いヘルパーを重ねるだけ）を踏襲しつつ、見た目の設計思想
だけを反転させたモジュールです。関数名も可能な限り pptx_skill.py と揃えて
あるので、既存コードの `from pptx_skill import ...` を
`from human_pptx import ...` に差し替えるだけでスタイルを切り替えられます。

pptx_skill.py との違い（詳細は SKILL.md 参照）:
- 背景は常に白 (`FFFFFF`)。アクセントバーやダーク背景の section_slide は無い
- タイトル文字は各スライド左上固定位置・18-20ptの見出しのみ（32pt+アクセント線を使わない）
- 箇条書きは "・" 始まりで、サブ項目も同じ記号でよい（ダッシュではなく中黒に統一）
- チャートはテーマの既定配色をそのまま使う（`series.format.fill` を一切触らない）
- 図解用に `add_box` / `add_arrow`（単色四角＋白文字、直線矢印）を追加
- 蛍光ペン風強調のための `highlight_run` を追加（python-pptx 標準APIには無いため OOXML を直接操作）
"""
from __future__ import annotations
from pathlib import Path
from typing import Sequence

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.oxml.ns import qn


# ---- Palette & fonts -------------------------------------------------------
# ここで"独自パレット"を設計しないのがポイント。白背景＋黒文字＋Office既定の
# チャート配色（テーマの accent1〜accent6 相当）をそのまま使う。

COLORS = {
    "bg":        RGBColor(0xFF, 0xFF, 0xFF),  # 常に白
    "text":      RGBColor(0x00, 0x00, 0x00),
    "sub":       RGBColor(0x59, 0x59, 0x59),
    "highlight": "FFFF00",                     # 蛍光ペン風（文字列でOOXMLにそのまま渡す）
    # 図解のボックス用（Office既定チャート配色の並びをそのまま流用）
    "box_blue":   RGBColor(0x44, 0x72, 0xC4),
    "box_orange": RGBColor(0xED, 0x7D, 0x31),
    "box_gray":   RGBColor(0xA5, 0xA5, 0xA5),
    "box_gold":   RGBColor(0xFF, 0xC0, 0x00),
    "box_cyan":   RGBColor(0x5B, 0x9B, 0xD5),
    "box_green":  RGBColor(0x70, 0xAD, 0x47),
    "white":      RGBColor(0xFF, 0xFF, 0xFF),
    "table_header_fill": RGBColor(0xD9, 0xE1, 0xF2),  # Excel既定っぽい薄い青
    "table_border": RGBColor(0xBF, 0xBF, 0xBF),
}

FONT_JP = "メイリオ"          # 見出し・本文の既定。QAはLibreOfficeで確認すること
FONT_JP_ALT = "游ゴシック"     # メイリオが無い環境向けの代替候補（明示切替用）
FONT_LATIN = "Arial"          # 英数字だけの場面用


# ---- Deck plumbing（pptx_skill.py と同一の考え方） -------------------------

def new_deck(widescreen: bool = True) -> Presentation:
    """16:9 既定（13.333"x7.5"）。4:3 にしたい場合は widescreen=False。"""
    prs = Presentation()
    if widescreen:
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
    else:
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
    return prs


def save(prs: Presentation, filename: str,
         out_dir: str = "/home/z/my-project/download") -> Path:
    p = Path(out_dir) / filename
    p.parent.mkdir(parents=True, exist_ok=True)
    prs.save(p)
    return p


def _blank(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _set_bg(slide, color: RGBColor = COLORS["bg"]) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _text(slide, left, top, width, height, text,
          size=14, bold=False, color=None,
          align=PP_ALIGN.LEFT, font=None, italic=False):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font or FONT_JP
    if color:
        run.font.color.rgb = color
    return box


# ---- 強調（蛍光ペン風ハイライト） ------------------------------------------

def highlight_run(run, hex_color: str = COLORS["highlight"]) -> None:
    """run のテキストに蛍光ペン風の背景色を付ける。

    python-pptx には run.font.highlight_color に相当する公開APIが無いため、
    OOXML の <a:highlight><a:srgbClr val="..."/></a:highlight> を直接差し込む。
    """
    rPr = run._r.get_or_add_rPr()
    # 既存の highlight があれば一旦削除してから追加（重複防止）
    existing = rPr.find(qn("a:highlight"))
    if existing is not None:
        rPr.remove(existing)
    highlight = rPr.makeelement(qn("a:highlight"), {})
    srgb = highlight.makeelement(qn("a:srgbClr"), {"val": hex_color})
    highlight.append(srgb)
    rPr.append(highlight)


# ---- 見出し（全スライド共通、左上固定） ------------------------------------

def add_header(slide, text: str, left=0.45, top=0.3, width=9.0, size=18):
    """各スライド左上の太字見出し。下線・色帯は付けない（アクセントライン厳禁）。"""
    return _text(slide, left, top, width, 0.6, text,
                 size=size, bold=True, color=COLORS["text"], font=FONT_JP)


# ---- 箇条書き（"・"始まり） -------------------------------------------------

def add_bullets(slide, left, top, width, height, items: Sequence[str],
                 size=14, sub_items: dict | None = None, mark="・"):
    """"・"始まりの箇条書き。sub_items={0: ["補足1", "補足2"]} でぶら下げ項目も置ける。

    因果関係を表す一行（例:「Aが多い　→　Bが顕著」）はitemsの文字列にそのまま
    "→" を書き込めばよい。矢印装飾を別途作る必要はない。
    """
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    first = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.text = f"{mark}{item}"
        p.font.size = Pt(size)
        p.font.name = FONT_JP
        p.font.color.rgb = COLORS["text"]
        p.space_after = Pt(8)
        if sub_items and i in sub_items:
            for sub in sub_items[i]:
                sp = tf.add_paragraph()
                sp.text = f"　　{mark}{sub}"
                sp.font.size = Pt(size - 2)
                sp.font.name = FONT_JP
                sp.font.color.rgb = COLORS["sub"]
                sp.space_after = Pt(4)
    return box


# ---- 図解（単色ボックス＋直線矢印） ----------------------------------------

def add_box(slide, left, top, width, height, label,
            fill_color: RGBColor = COLORS["box_blue"],
            font_color: RGBColor = COLORS["white"], size=14, shape=MSO_SHAPE.ROUNDED_RECTANGLE):
    """図解用の単色ボックス。影・グラデーションは付けない。"""
    box = slide.shapes.add_shape(shape, Inches(left), Inches(top), Inches(width), Inches(height))
    box.fill.solid()
    box.fill.fore_color.rgb = fill_color
    box.line.color.rgb = fill_color
    box.shadow.inherit = False
    # add_shape() attaches a theme <p:style> with an effectRef (=drop shadow) by
    # default; an empty <a:effectLst/> alone doesn't reliably suppress it in
    # every renderer, so drop the whole <p:style> element to be sure.
    style_el = box._element.find(qn("p:style"))
    if style_el is not None:
        box._element.remove(style_el)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = label
    run.font.size = Pt(size)
    run.font.bold = True
    run.font.name = FONT_JP
    run.font.color.rgb = font_color
    return box


def add_arrow(slide, x1, y1, x2, y2, color: RGBColor = COLORS["text"], weight_pt=1.5):
    """シンプルな直線矢印（曲線・アイコン装飾なし）。座標はすべてインチ。"""
    connector = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )
    connector.line.color.rgb = color
    connector.line.width = Pt(weight_pt)
    # 矢印の先端（終点側）を付ける
    ln = connector.line._get_or_add_ln()
    tail = ln.makeelement(qn("a:tailEnd"), {"type": "triangle"})
    ln.append(tail)
    return connector


def add_photo(slide, image_path, left, top, width=None, height=None):
    """実写真・イラストをそのまま貼るだけ。角丸マスクや枠線は付けない。"""
    kwargs = {}
    if width is not None:
        kwargs["width"] = Inches(width)
    if height is not None:
        kwargs["height"] = Inches(height)
    return slide.shapes.add_picture(str(image_path), Inches(left), Inches(top), **kwargs)


# ---- スライド種別（pptx_skill.py と対になる関数名） -------------------------

def cover_slide(prs: Presentation, title: str, subtitle: str = "", authors: str = ""):
    """表紙。白背景、中央揃え、装飾なし。"""
    s = _blank(prs)
    _set_bg(s)
    _text(s, 1.0, 2.6, 11.3, 1.2, title, size=30, bold=True,
          color=COLORS["text"], align=PP_ALIGN.CENTER, font=FONT_JP)
    if subtitle:
        _text(s, 1.0, 3.7, 11.3, 0.8, subtitle, size=17,
              color=COLORS["sub"], align=PP_ALIGN.CENTER, font=FONT_JP)
    if authors:
        _text(s, 1.0, 6.5, 11.3, 0.6, authors, size=13,
              color=COLORS["sub"], align=PP_ALIGN.CENTER, font=FONT_JP)
    return s


def bullet_slide(prs: Presentation, header: str, bullets: Sequence[str],
                  sub_items: dict | None = None, size=14):
    s = _blank(prs)
    _set_bg(s)
    add_header(s, header)
    add_bullets(s, 0.5, 1.15, 11.5, 5.5, bullets, size=size, sub_items=sub_items)
    return s


def image_text_slide(prs: Presentation, header: str, image_path, text_bullets: Sequence[str],
                      image_side: str = "left", size=14):
    s = _blank(prs)
    _set_bg(s)
    add_header(s, header)
    if image_side == "left":
        add_photo(s, image_path, 0.5, 1.2, width=5.6)
        add_bullets(s, 6.4, 1.3, 6.0, 5.3, text_bullets, size=size)
    else:
        add_bullets(s, 0.5, 1.3, 6.0, 5.3, text_bullets, size=size)
        add_photo(s, image_path, 6.9, 1.2, width=5.6)
    return s


def chart_slide(prs: Presentation, header: str, categories: Sequence[str],
                 series_data: dict, chart_type=XL_CHART_TYPE.COLUMN_CLUSTERED,
                 note: str = ""):
    """ネイティブチャート。色・スタイルは一切カスタマイズせず、テーマの既定配色に任せる。"""
    s = _blank(prs)
    _set_bg(s)
    add_header(s, header)

    cd = CategoryChartData()
    cd.categories = list(categories)
    for name, values in series_data.items():
        cd.add_series(name, list(values))

    gframe = s.shapes.add_chart(
        chart_type, Inches(0.8), Inches(1.2), Inches(11.5), Inches(5.3), cd
    )
    chart = gframe.chart
    chart.has_legend = len(series_data) > 1
    chart.has_title = False
    # フォントだけ日本語に合わせる。配色・データラベル位置などは既定のまま触らない
    if chart.has_legend:
        chart.legend.font.size = Pt(11)
        chart.legend.font.name = FONT_JP
    for axis in (getattr(chart, "category_axis", None), getattr(chart, "value_axis", None)):
        if axis is not None:
            axis.tick_labels.font.size = Pt(10)
            axis.tick_labels.font.name = FONT_JP

    if note:
        _text(s, 0.8, 6.75, 11.5, 0.4, note, size=9, color=COLORS["sub"], font=FONT_JP)
    return s


def table_slide(prs: Presentation, header: str, headers: Sequence[str],
                 rows: Sequence[Sequence[object]], size=13):
    s = _blank(prs)
    _set_bg(s)
    add_header(s, header)

    rows_count = len(rows) + 1
    cols_count = len(headers)
    tbl_shape = s.shapes.add_table(
        rows_count, cols_count, Inches(0.5), Inches(1.2), Inches(11.5), Inches(5.2)
    )
    tbl = tbl_shape.table

    for i, h in enumerate(headers):
        cell = tbl.cell(0, i)
        cell.text = str(h)
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLORS["table_header_fill"]
        for p in cell.text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER
            for r in p.runs:
                r.font.color.rgb = COLORS["text"]
                r.font.bold = True
                r.font.size = Pt(size)
                r.font.name = FONT_JP

    for ri, row in enumerate(rows, start=1):
        for ci, v in enumerate(row):
            cell = tbl.cell(ri, ci)
            cell.text = str(v)
            for p in cell.text_frame.paragraphs:
                for r in p.runs:
                    r.font.size = Pt(size)
                    r.font.name = FONT_JP
    return s


def diagram_slide(prs: Presentation, header: str, boxes: Sequence[dict],
                   arrows: Sequence[tuple] | None = None):
    """システム構成図・フロー図用のスライド。

    boxes: [{"label": "入力", "left":.., "top":.., "width":.., "height":..,
             "color": COLORS["box_blue"]}, ...]
    arrows: [(x1, y1, x2, y2), ...]  座標はすべてインチ
    """
    s = _blank(prs)
    _set_bg(s)
    add_header(s, header)
    for b in boxes:
        add_box(
            s, b["left"], b["top"], b["width"], b["height"], b["label"],
            fill_color=b.get("color", COLORS["box_blue"]),
            font_color=b.get("font_color", COLORS["white"]),
            size=b.get("size", 14),
        )
    for a in arrows or []:
        add_arrow(s, *a)
    return s


# ---- Self-test --------------------------------------------------------------

if __name__ == "__main__":
    # 動作確認用のダミーデータ（架空の題材）。実プロジェクトの内容には一切依存しない。
    prs = new_deck()
    cover_slide(prs, "サンプル研究発表のタイトル",
                "〜サブタイトルはここに入れる〜", "発表者名A　発表者名B　発表者名C")
    bullet_slide(prs, "背景・目的", [
        "課題Aと課題Bには関連がある",
        "指標Xは目標値Yに届いていない　→　特に条件Zで差が大きい",
        "本研究では条件に応じた提案システムを構築し、課題の解決を図る",
    ])
    chart_slide(prs, "背景・目的", ["区分1", "区分2", "区分3", "区分4", "区分5", "区分6"],
                {"指標X": [24.1, 25.6, 26.0, 27.0, 28.0, 29.0]},
                note="出典：サンプル出典（実際の発表では出典を明記する）")
    diagram_slide(prs, "システム構成 - 完成図", boxes=[
        {"label": "入力", "left": 0.8, "top": 2.5, "width": 2.2, "height": 1.0, "color": COLORS["box_blue"]},
        {"label": "検出処理", "left": 4.0, "top": 2.5, "width": 2.2, "height": 1.0, "color": COLORS["box_cyan"]},
        {"label": "生成処理", "left": 7.2, "top": 2.5, "width": 2.2, "height": 1.0, "color": COLORS["box_gold"]},
        {"label": "利用者", "left": 10.4, "top": 2.5, "width": 2.2, "height": 1.0, "color": COLORS["box_green"]},
    ], arrows=[(3.0, 3.0, 4.0, 3.0), (6.2, 3.0, 7.2, 3.0), (9.4, 3.0, 10.4, 3.0)])
    out = save(prs, "human_pptx_demo.pptx")
    print(f"Saved: {out}")
